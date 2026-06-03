"""
圣经电影全自动生成工作台 — /film-studio
依赖: anthropic google-genai edge-tts boto3 pillow (ffmpeg 系统包)
环境变量: ANTHROPIC_API_KEY, GEMINI_API_KEY,
          R2_ACCOUNT_ID, R2_ACCESS_KEY_ID, R2_SECRET_ACCESS_KEY,
          R2_BUCKET_NAME, VIDEO_CDN_BASE
"""

import os, re, sys, json, time, uuid, asyncio, threading, subprocess, io, textwrap
from pathlib import Path
from typing import Generator
import httpx
from fastapi import APIRouter, BackgroundTasks, UploadFile, File, Form
from fastapi.responses import HTMLResponse, FileResponse, StreamingResponse
from pydantic import BaseModel

# ── 路由 & 存储目录 ────────────────────────────────────────────────────────────
router = APIRouter(tags=["film-studio"])

FILM_DIR  = Path("/app/film_output")      # HF Docker 路径
if not Path("/app").exists():
    FILM_DIR = Path("./film_output")       # 本地开发
CLIPS_DIR  = FILM_DIR / "clips"
AUDIO_DIR  = FILM_DIR / "audio"
COMP_DIR   = FILM_DIR / "composed"
SLIDES_DIR = FILM_DIR / "slides"
UPLOAD_DIR = FILM_DIR / "uploads"
for _d in [CLIPS_DIR, AUDIO_DIR, COMP_DIR, SLIDES_DIR, UPLOAD_DIR]:
    _d.mkdir(parents=True, exist_ok=True)

JOBS: dict[str, dict] = {}
FONT_PATH = next(
    (p for p in [
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ] if Path(p).exists()),
    ""
)


# ══════════════════════════════════════════════════════════════════════════════
# 流水线函数
# ══════════════════════════════════════════════════════════════════════════════

def _ff(cmd: list) -> bool:
    r = subprocess.run(cmd, capture_output=True)
    if r.returncode != 0:
        print(f"[FFmpeg] {r.stderr[-500:]}")
    return r.returncode == 0


def split_with_claude(story_text: str, api_key: str, n: int) -> dict:
    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    system = textwrap.dedent(f"""
        You are a biblical film director. Given a storyboard, produce exactly {n} scene entries
        plus 1 spiritual_application block. Return ONLY valid JSON, no markdown, in this schema:
        {{
          "title": "film title",
          "scenes": [{{
            "id": 1,
            "video_prompt": "Detailed Veo 3.1 prompt (3-5 sentences). Camera angle, lighting, character, action, emotion. End with: 16:9 aspect ratio, 4K cinematic, no text, no subtitles.",
            "narration_zh": "Chinese narration 15-30 chars (spoken aloud during clip).",
            "subtitle_zh": "Chinese subtitle 8-18 chars (shown at screen bottom)."
          }}],
          "spiritual_application": {{
            "title_zh": "结语标题 (e.g. 约瑟的故事告诉我们：)",
            "lines_zh": ["第一行", "第二行", "第三行"],
            "scripture_zh": "经文引用",
            "narration_zh": "完整旁白 60-100字",
            "duration_sec": 20
          }}
        }}
        ALL narration/subtitle text must be Simplified Chinese (简体中文).
        Keep character descriptions IDENTICAL across all scene prompts.
    """)
    resp = client.messages.create(
        model="claude-opus-4-5", max_tokens=8192, system=system,
        messages=[{"role": "user", "content": f"Storyboard:\n\n{story_text}"}],
    )
    raw = re.sub(r'^```[a-z]*\n?', '', resp.content[0].text.strip())
    raw = re.sub(r'\n?```$', '', raw)
    return json.loads(raw)



def _split_system_prompt(n: int) -> str:
    return textwrap.dedent(f"""
        You are a biblical film director. Given a storyboard, produce exactly {n} scene entries
        plus 1 spiritual_application block. Return ONLY valid JSON, no markdown, in this schema:
        {{
          "title": "film title",
          "scenes": [{{
            "id": 1,
            "video_prompt": "Detailed Veo 3.1 prompt (3-5 sentences). Camera angle, lighting, character, action, emotion. End with: 16:9 aspect ratio, 4K cinematic, no text, no subtitles.",
            "narration_zh": "Chinese narration 15-30 chars (spoken aloud during clip).",
            "subtitle_zh": "Chinese subtitle 8-18 chars (shown at screen bottom)."
          }}],
          "spiritual_application": {{
            "title_zh": "结语标题 (e.g. 约瑟的故事告诉我们：)",
            "lines_zh": ["第一行", "第二行", "第三行"],
            "scripture_zh": "经文引用",
            "narration_zh": "完整旁白 60-100字",
            "duration_sec": 20
          }}
        }}
        ALL narration/subtitle text must be Simplified Chinese (简体中文).
        Keep character descriptions IDENTICAL across all scene prompts.
    """)


def split_with_gemini(story_text: str, api_key: str, n: int) -> dict:
    """Split a storyboard into scenes using Gemini (no Anthropic credits needed)."""
    from google import genai
    from google.genai import types
    client = genai.Client(api_key=api_key)
    model = os.environ.get("GEMINI_TEXT_MODEL", "gemini-2.5-flash")
    resp = client.models.generate_content(
        model=model,
        contents=f"Storyboard:\n\n{story_text}",
        config=types.GenerateContentConfig(
            system_instruction=_split_system_prompt(n),
            response_mime_type="application/json",
            temperature=0.7,
            max_output_tokens=32768,
            thinking_config=types.ThinkingConfig(thinking_budget=0),
        ),
    )
    raw = (resp.text or "").strip()
    raw = re.sub(r'^```[a-z]*\n?', '', raw)
    raw = re.sub(r'\n?```$', '', raw)
    return json.loads(raw)


class SpendCapExceeded(Exception):
    """Veo/Gemini 项目月度支出上限已满（429 RESOURCE_EXHAUSTED）。"""


def _is_spend_cap(e) -> bool:
    msg = str(e).lower()
    return ("resource_exhausted" in msg or "spend cap" in msg
            or "spending cap" in msg or "'code': 429" in msg or "code: 429" in msg)


def generate_veo_clip(prompt: str, path: Path, api_key: str, cb=None, fallback_key: str = "") -> bool:
    from google import genai
    from google.genai import types
    # 依次尝试主 key 与备用 key（去重）；某把 key 被 referrer/权限拦截时自动换下一把
    keys = list(dict.fromkeys([k for k in (api_key, fallback_key) if k]))
    last_err = None
    spend_cap = False
    for ki, key in enumerate(keys):
        try:
            client = genai.Client(api_key=key)
            op = client.models.generate_videos(
                model="veo-3.1-generate-preview", prompt=prompt,
                config=types.GenerateVideosConfig(aspect_ratio="16:9"),
            )
            waited = 0
            while not op.done:
                time.sleep(15); waited += 15
                op = client.operations.get(op)
                if cb: cb(f"Veo {waited}s…")
                if waited >= 660: raise TimeoutError("Veo 超时")
            uri = op.result.generated_videos[0].video.uri
            url = uri + (f"&key={key}" if "googleapis.com" in uri and "key=" not in uri else "")
            with httpx.Client(timeout=120, follow_redirects=True) as hc:
                with hc.stream("GET", url) as r:
                    r.raise_for_status()
                    with open(path, "wb") as f:
                        for chunk in r.iter_bytes(65536): f.write(chunk)
            return True
        except Exception as e:
            last_err = e
            if _is_spend_cap(e): spend_cap = True
            print(f"[Veo] key#{ki+1}/{len(keys)} {e}")
            if cb and ki + 1 < len(keys): cb(f"key#{ki+1} 失败，换备用 key 重试…")
            continue
    if spend_cap:
        raise SpendCapExceeded("Gemini 项目月度支出上限已满（spend cap）——请到 https://ai.studio/spend 调高后重试")
    print(f"[Veo] 所有 key 均失败: {last_err}")
    return False


def _tts_elevenlabs(text: str, path: Path) -> bool:
    """ElevenLabs 高质量配音（配置 ELEVENLABS_API_KEY 时优先）。"""
    key = os.environ.get("ELEVENLABS_API_KEY", "")
    if not key:
        return False
    try:
        voice = os.environ.get("ELEVENLABS_VOICE_ID", "21m00Tcm4TlvDq8ikWAM")
        model = os.environ.get("ELEVENLABS_MODEL", "eleven_multilingual_v2")
        url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice}"
        with httpx.Client(timeout=120) as hc:
            r = hc.post(url, headers={"xi-api-key": key, "accept": "audio/mpeg",
                                      "content-type": "application/json"},
                        json={"text": text, "model_id": model,
                              "voice_settings": {"stability": 0.4, "similarity_boost": 0.8}})
            r.raise_for_status()
            path.write_bytes(r.content)
        if path.exists() and path.stat().st_size > 256:
            print("[TTS] via ElevenLabs", flush=True); return True
    except Exception as e:
        print(f"[TTS] ElevenLabs failed: {e}", flush=True)
    return False


def _tts_gtts(text: str, path: Path) -> bool:
    """谷歌在线 TTS 兜底（与微软不同网络路径）。"""
    try:
        from gtts import gTTS
        gTTS(text=text, lang="zh-CN").save(str(path))
        if path.exists() and path.stat().st_size > 256:
            print("[TTS] via gTTS", flush=True); return True
    except Exception as e:
        print(f"[TTS] gTTS failed: {e}", flush=True)
    return False


def _tts_espeak(text: str, path: Path) -> bool:
    """espeak-ng 离线 TTS 兜底（音质机械但无需联网，保证一定有声）。"""
    try:
        wav = path.with_suffix(".espeak.wav")
        subprocess.run(["espeak-ng", "-v", "cmn", "-s", "150", "-w", str(wav), text],
                       capture_output=True, timeout=60)
        if wav.exists() and wav.stat().st_size > 256:
            ok = _ff(["ffmpeg", "-y", "-i", str(wav), "-acodec", "libmp3lame", "-q:a", "4", str(path)])
            try: wav.unlink()
            except Exception: pass
            if ok and path.exists() and path.stat().st_size > 256:
                print("[TTS] via espeak-ng (offline)", flush=True); return True
    except Exception as e:
        print(f"[TTS] espeak-ng failed: {e}", flush=True)
    return False


async def tts_to_file(text: str, path: Path, use_elevenlabs: bool = True) -> bool:
    """多级兜底：edge-tts(微软) → gTTS(谷歌) → espeak-ng(离线) → 静音占位。"""
    text = (text or "").strip()
    if text:
        # 0) ElevenLabs —— 最高音质（在线，需 ELEVENLABS_API_KEY 且页面勾选）
        if use_elevenlabs and _tts_elevenlabs(text, path): return True
        # 1) edge-tts —— 音质良好（在线·微软）
        try:
            import edge_tts
            buf = io.BytesIO()
            async for chunk in edge_tts.Communicate(text, "zh-CN-XiaoxiaoNeural", rate="+0%").stream():
                if chunk["type"] == "audio": buf.write(chunk["data"])
            if buf.tell() > 0:
                path.write_bytes(buf.getvalue())
                print("[TTS] via edge-tts", flush=True); return True
        except Exception as e:
            print(f"[TTS] edge-tts failed: {e}", flush=True)
        # 2) gTTS（在线·谷歌）
        if _tts_gtts(text, path): return True
        # 3) espeak-ng（离线，保证有声）
        if _tts_espeak(text, path): return True
    # 4) 静音占位（无文本或全部失败）
    _ff(["ffmpeg", "-y", "-f", "lavfi", "-i", "anullsrc=r=24000:cl=mono",
         "-t", "7", "-q:a", "9", "-acodec", "libmp3lame", str(path)])
    return False


def compose_clip(video: Path, audio: Path, subtitle: str, out: Path) -> bool:
    fo = f":fontfile={FONT_PATH}" if FONT_PATH else ""
    esc = subtitle.replace("'", "\\'").replace(":", "\\:")
    dt = (f"drawtext=text='{esc}'{fo}:fontsize=34:fontcolor=white"
          f":x=(w-text_w)/2:y=h-75:borderw=3:bordercolor=black@0.8"
          f":shadowx=2:shadowy=2:shadowcolor=black@0.5")
    return _ff([
        "ffmpeg","-y","-i",str(video),"-i",str(audio),
        "-vf",dt,"-c:v","libx264","-preset","fast","-crf","22",
        "-c:a","aac","-b:a","128k","-shortest","-pix_fmt","yuv420p",str(out),
    ])


def create_spiritual_scene(sp: dict, audio: Path, out: Path) -> bool:
    title  = sp.get("title_zh","属灵应用")
    lines  = sp.get("lines_zh",[])
    scr    = sp.get("scripture_zh","")
    dur    = sp.get("duration_sec",20)
    bg     = FILM_DIR / "sp_bg.mp4"

    try:
        from PIL import Image, ImageDraw, ImageFont
        W,H = 1920,1080
        img = Image.new("RGB",(W,H),(10,10,26))
        draw = ImageDraw.Draw(img)
        def font(sz):
            try: return ImageFont.truetype(FONT_PATH,sz) if FONT_PATH else ImageFont.load_default()
            except: return ImageFont.load_default()
        y = 220
        draw.text((W//2,y), title, font=font(52), fill=(255,214,10), anchor="mm")
        y += 90
        draw.line([(W//2-180,y),(W//2+180,y)], fill=(255,214,10,120), width=2)
        y += 50
        for ln in lines:
            draw.text((W//2,y), ln, font=font(40), fill=(235,235,255), anchor="mm"); y += 72
        if scr:
            draw.text((W//2,y+20), f"— {scr}", font=font(30), fill=(170,150,255), anchor="mm")
        png = FILM_DIR/"sp_bg.png"; img.save(str(png))
        _ff(["ffmpeg","-y","-loop","1","-i",str(png),"-t",str(dur),
             "-c:v","libx264","-pix_fmt","yuv420p","-vf","scale=1920:1080",str(bg)])
    except Exception as e:
        print(f"[Spiritual PIL] {e}")
        fo = f":fontfile={FONT_PATH}" if FONT_PATH else ""
        esc = (title+" "+" ".join(lines)).replace("'","\\'").replace(":",r"\:")
        _ff(["ffmpeg","-y","-f","lavfi",
             "-i",f"color=c=0x0a0a1a:size=1920x1080:rate=25","-t",str(dur),
             "-vf",f"drawtext=text='{esc}'{fo}:fontsize=38:fontcolor=white:x=(w-text_w)/2:y=(h-text_h)/2:borderw=3:bordercolor=black",
             "-c:v","libx264","-pix_fmt","yuv420p",str(bg)])

    return _ff([
        "ffmpeg","-y","-i",str(bg),"-i",str(audio),
        "-c:v","copy","-c:a","aac","-b:a","128k","-shortest",str(out),
    ])


_NORM_VF = ("scale=1920:1080:force_original_aspect_ratio=decrease,"
            "pad=1920:1080:(ow-iw)/2:(oh-ih)/2,setsar=1,fps=30,format=yuv420p")


def _has_audio(p: Path) -> bool:
    r = subprocess.run(
        ["ffprobe","-v","error","-select_streams","a","-show_entries","stream=index",
         "-of","csv=p=0",str(p)], capture_output=True, text=True)
    return bool(r.stdout.strip())


def _normalize_clip(src: Path, dst: Path) -> bool:
    """统一分辨率/帧率/像素格式/音轨（无音轨补静音），便于 concat。"""
    if _has_audio(src):
        cmd = ["ffmpeg","-y","-i",str(src),"-vf",_NORM_VF,
               "-c:v","libx264","-preset","fast","-crf","20",
               "-c:a","aac","-ar","48000","-ac","2","-b:a","192k",
               "-r","30",str(dst)]
    else:
        cmd = ["ffmpeg","-y","-i",str(src),
               "-f","lavfi","-i","anullsrc=channel_layout=stereo:sample_rate=48000",
               "-vf",_NORM_VF,"-map","0:v:0","-map","1:a:0",
               "-c:v","libx264","-preset","fast","-crf","20",
               "-c:a","aac","-ar","48000","-ac","2","-b:a","192k",
               "-r","30","-shortest",str(dst)]
    return _ff(cmd)


def concat_all(clips: list[Path], out: Path) -> bool:
    norm_dir = FILM_DIR/"norm"; norm_dir.mkdir(parents=True, exist_ok=True)
    normed: list[Path] = []
    for i, p in enumerate(clips):
        if not (p.exists() and p.stat().st_size > 1024):
            continue
        d = norm_dir / f"n_{i:03d}.mp4"
        if _normalize_clip(p, d) and d.exists() and d.stat().st_size > 1024:
            normed.append(d)
        else:
            print(f"[FFmpeg] 归一化失败，跳过 {p.name}")
    if not normed:
        print("[FFmpeg] 无有效片段可拼接"); return False
    lst = FILM_DIR/"concat.txt"
    with open(lst,"w") as f:
        for d in normed:
            f.write(f"file '{d.resolve()}'\n")
    return _ff([
        "ffmpeg","-y","-f","concat","-safe","0","-i",str(lst),
        "-c:v","libx264","-preset","medium","-crf","20",
        "-c:a","aac","-b:a","192k","-pix_fmt","yuv420p","-movflags","+faststart",str(out),
    ])


def upload_r2(path: Path, prefix: str = "biblical-films/") -> str:
    import boto3
    aid = os.environ.get("R2_ACCOUNT_ID","")
    ak  = os.environ.get("R2_ACCESS_KEY_ID","")
    sk  = os.environ.get("R2_SECRET_ACCESS_KEY","")
    bkt = os.environ.get("R2_BUCKET_NAME","")
    cdn = os.environ.get("VIDEO_CDN_BASE", f"https://{bkt}.r2.dev").rstrip("/")
    if not all([aid,ak,sk,bkt]): raise ValueError("R2 env vars missing")
    s3 = boto3.client("s3", endpoint_url=f"https://{aid}.r2.cloudflarestorage.com",
                      aws_access_key_id=ak, aws_secret_access_key=sk, region_name="auto")
    key = prefix + path.name
    s3.upload_file(str(path), bkt, key, ExtraArgs={"ContentType":"video/mp4"})
    return f"{cdn}/{key}"


# ══════════════════════════════════════════════════════════════════════════════
# 主流水线（后台线程）
# ══════════════════════════════════════════════════════════════════════════════

def _log(job, msg, pct=None):
    job["steps"].append(msg)
    if pct is not None: job["progress"] = pct
    print(f"[Film] {msg}")


def run_pipeline(job_id: str, story: str, ak: str, gk: str, ck: str, n: int):
    job = JOBS[job_id]
    job.update(status="running", steps=[], progress=0)
    fid = job_id[:8]

    try:
        # Step 1: Gemini 拆分（不消耗 Anthropic 额度）
        _log(job, "🤖 Gemini 拆分镜头…", 3)
        data   = split_with_gemini(story, ck, n)
        scenes = data["scenes"]
        sp     = data.get("spiritual_application", {})
        job["story"] = data
        _log(job, f"✅ 共 {len(scenes)} 个镜头", 8)

        composed: list[Path] = []

        # Step 2-4: 逐镜头
        for i, sc in enumerate(scenes):
            sid  = sc["id"]
            base = 8 + int(i/len(scenes)*75)
            _log(job, f"🎬 Scene {sid}/{len(scenes)}: Veo 生成…", base)
            job["cur"] = sid

            clip  = CLIPS_DIR / f"{fid}_{sid:02d}.mp4"
            aud   = AUDIO_DIR / f"{fid}_{sid:02d}.mp3"
            comp  = COMP_DIR  / f"{fid}_{sid:02d}.mp4"

            if not (clip.exists() and clip.stat().st_size > 1024):
                ok = generate_veo_clip(sc["video_prompt"], clip, ck,
                                       cb=lambda m: _log(job, f"  ·{m}"), fallback_key=gk)
                if not ok:
                    _log(job, f"  ⚠️ Scene {sid} Veo 失败，跳过（不计入拼接）")
                    time.sleep(3); continue
            else:
                _log(job, f"  ↩ Scene {sid} 复用已有片段")

            asyncio.run(tts_to_file(sc.get("narration_zh",""), aud))
            compose_clip(clip, aud, sc.get("subtitle_zh",""), comp)
            composed.append(comp if (comp.exists() and comp.stat().st_size>1024) else clip)
            _log(job, f"  ✅ Scene {sid} 完成", base+2)
            time.sleep(3)

        # Step 5: 属灵应用结尾
        _log(job, "✨ 属灵应用结尾…", 85)
        sp_aud = AUDIO_DIR / f"{fid}_sp.mp3"
        sp_vid = COMP_DIR  / f"{fid}_sp.mp4"
        asyncio.run(tts_to_file(sp.get("narration_zh",""), sp_aud))
        create_spiritual_scene(sp, sp_aud, sp_vid)
        if sp_vid.exists(): composed.append(sp_vid)

        # Step 6: 拼接
        _log(job, "🔗 FFmpeg 拼接…", 90)
        final = FILM_DIR / f"{fid}_final.mp4"
        if not concat_all(composed, final):
            raise RuntimeError("FFmpeg 拼接失败")
        mb = round(final.stat().st_size/1024/1024, 1)
        _log(job, f"✅ 拼接完成 {mb} MB", 95)

        # Step 7: R2 上传
        r2_url = None
        try:
            _log(job, "☁️ 上传 R2…", 97)
            r2_url = upload_r2(final)
            _log(job, f"✅ {r2_url}", 99)
        except Exception as e:
            _log(job, f"⚠️ R2 跳过: {e}")

        job.update(status="done", progress=100,
                   result={"file": final.name, "r2_url": r2_url, "mb": mb, "scenes": len(scenes)})
        _log(job, "🎉 完成！")

    except Exception as e:
        import traceback; traceback.print_exc()
        job.update(status="error", error=str(e))
        _log(job, f"❌ {e}")


# ── PPT → 视频管线（PPT插画 + Ken Burns + edge-tts，离线无 LLM）─────────────────

def pptx_to_images(pptx_path: Path, out_dir: Path) -> list[Path]:
    """用 LibreOffice 把 pptx 转 pdf，再用 poppler 逐页转 png。"""
    out_dir.mkdir(parents=True, exist_ok=True)
    prof = f"file:///tmp/lo_{out_dir.name}"   # 独立 profile，避免并发锁
    subprocess.run(
        ["soffice", "--headless", f"-env:UserInstallation={prof}",
         "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
        capture_output=True, timeout=240)
    pdf = out_dir / (pptx_path.stem + ".pdf")
    if not pdf.exists():
        cands = list(out_dir.glob("*.pdf"))
        if not cands: raise RuntimeError("PPT 转 PDF 失败（LibreOffice 未产出 pdf）")
        pdf = cands[0]
    subprocess.run(["pdftoppm", "-r", "150", "-png", str(pdf), str(out_dir / "slide")],
                   capture_output=True, timeout=240)
    imgs = list(out_dir.glob("slide-*.png"))
    imgs.sort(key=lambda q: int(re.search(r"(\d+)", q.stem).group(1)))
    return imgs


def pptx_notes(pptx_path: Path) -> list[str]:
    """逐页读取演讲者备注作为旁白文本。"""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    out = []
    for sl in prs.slides:
        txt = ""
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame is not None:
            txt = (sl.notes_slide.notes_text_frame.text or "").strip()
        out.append(txt)
    return out


def pptx_slide_texts(pptx_path: Path) -> list[str]:
    """逐页读取幻灯片上的文字（备注为空时作为旁白回退来源）。"""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    out = []
    for sl in prs.slides:
        parts = []
        for shp in sl.shapes:
            if shp.has_text_frame:
                t = (shp.text_frame.text or "").strip()
                if t:
                    parts.append(t)
        out.append("\n".join(parts).strip())
    return out


def _audio_dur(p: Path) -> float:
    try:
        r = subprocess.run(["ffprobe", "-v", "error", "-show_entries", "format=duration",
                            "-of", "csv=p=0", str(p)], capture_output=True, text=True)
        return float(r.stdout.strip())
    except Exception:
        return 0.0


def kenburns_clip(img: Path, dur: float, out: Path, zoom_in: bool = True) -> bool:
    """给单张图加缓慢推近/拉远的镜头运动，输出 1920x1080 无声片段。"""
    frames = max(1, int(dur * 30))
    z = "min(zoom+0.0008,1.18)" if zoom_in else "if(lte(zoom,1.0),1.18,max(1.001,zoom-0.0008))"
    vf = ("scale=2400:1350:force_original_aspect_ratio=increase,crop=2400:1350,"
          "zoompan=z='" + z + "':x='iw/2-(iw/zoom/2)':y='ih/2-(ih/zoom/2)':"
          "d=" + str(frames) + ":s=1920x1080:fps=30,setsar=1,format=yuv420p")
    return _ff(["ffmpeg", "-y", "-loop", "1", "-i", str(img), "-t", f"{dur:.2f}", "-r", "30",
                "-filter_complex", vf, "-c:v", "libx264", "-preset", "fast", "-crf", "20", "-an", str(out)])


def _kling_jwt(ak: str, sk: str) -> str:
    """官方算法：JWT(HS256) payload={iss:ak, exp:+1800s, nbf:-5s}。
    注意 PyJWT 1.x 返回 bytes，直接拼 Bearer 会变成 b'...' 导致 401，这里统一转 str。"""
    import jwt as _jwt
    now = int(time.time())
    tok = _jwt.encode({"iss": ak, "exp": now + 1800, "nbf": now - 5}, sk,
                      algorithm="HS256", headers={"alg": "HS256", "typ": "JWT"})
    return tok.decode() if isinstance(tok, bytes) else tok


def kling_configured() -> bool:
    return bool(os.environ.get("KLING_ACCESS_KEY") and os.environ.get("KLING_SECRET_KEY"))


_KLING_BASE_OK: "str | None" = None  # 记住验证通过的端点，避免每页重复探测

def _kling_bases() -> list:
    """候选端点：env 指定的优先，其后按区域兜底。Kling 的 AK/SK 分区域，
    国际版 key 只能打 api-singapore，国内开放平台 key 只能打 api-beijing，打错=401。"""
    env = os.environ.get("KLING_API_BASE", "").rstrip("/")
    bases = [env] if env else []
    for b in ("https://api-singapore.klingai.com",
              "https://api-beijing.klingai.com",
              "https://api.klingai.com"):
        if b not in bases:
            bases.append(b)
    return bases


def _pad_video(video: Path, pad_sec: float, out: Path) -> bool:
    """补帧(克隆末帧)把视频延长 pad_sec 秒，用于与较长旁白对齐。"""
    if pad_sec <= 0.05:
        return False
    return _ff(["ffmpeg", "-y", "-i", str(video),
                "-vf", f"tpad=stop_mode=clone:stop_duration={pad_sec:.2f},format=yuv420p",
                "-c:v", "libx264", "-preset", "fast", "-crf", "20", "-an", str(out)])


def generate_kling_clip(image: Path, prompt: str, dur_sec: float, out: Path, cb=None) -> bool:
    """Kling 图生视频：把单张插画变成真实动画。失败返回 False（上层回退 Ken Burns）。"""
    ak = os.environ.get("KLING_ACCESS_KEY", "")
    sk = os.environ.get("KLING_SECRET_KEY", "")
    if not (ak and sk):
        return False
    try:
        import base64
        global _KLING_BASE_OK
        model = os.environ.get("KLING_MODEL", "kling-v1")
        mode = os.environ.get("KLING_MODE", "std")        # std / pro
        kdur = "10" if dur_sec > 5.5 else "5"
        b64 = base64.b64encode(image.read_bytes()).decode()
        body = {"model_name": model, "image": b64, "mode": mode, "duration": kdur,
                "prompt": (prompt or "圣经故事场景，电影感，自然真实的人物与环境动作")[:2000]}
        with httpx.Client(timeout=60, limits=httpx.Limits(max_keepalive_connections=0)) as hc:
            r = None
            base = None
            candidates = [_KLING_BASE_OK] if _KLING_BASE_OK else _kling_bases()
            for cand in candidates:
                try:
                    resp = hc.post(f"{cand}/v1/videos/image2video",
                                   headers={"Authorization": f"Bearer {_kling_jwt(ak, sk)}",
                                            "Content-Type": "application/json"}, json=body)
                except Exception as ce:
                    print(f"[Kling] 创建任务网络错误 @ {cand}: {ce} — 重试一次", flush=True)
                    try:
                        resp = hc.post(f"{cand}/v1/videos/image2video",
                                       headers={"Authorization": f"Bearer {_kling_jwt(ak, sk)}",
                                                "Content-Type": "application/json"}, json=body)
                    except Exception as ce2:
                        print(f"[Kling] 仍失败 @ {cand}: {ce2}", flush=True)
                        continue
                if resp.status_code == 401:
                    print(f"[Kling] 401 @ {cand} — AK/SK 与该区域端点不匹配，试下一个端点", flush=True)
                    continue
                resp.raise_for_status()
                r = resp
                base = cand
                _KLING_BASE_OK = cand
                print(f"[Kling] 端点可用: {cand}", flush=True)
                break
            if r is None:
                print("[Kling] 所有端点均 401：请确认 AK/SK 来自哪个平台（国际版→api-singapore / 国内可灵开放平台→api-beijing），且账号已开通 API 资源包", flush=True)
                return False
            tid = (r.json().get("data") or {}).get("task_id")
            if not tid:
                print(f"[Kling] no task_id: {r.text[:200]}", flush=True); return False
            waited = 0
            poll_fail = 0
            while waited < 600:
                time.sleep(10); waited += 10
                try:
                    sresp = hc.get(f"{base}/v1/videos/image2video/{tid}",
                                   headers={"Authorization": f"Bearer {_kling_jwt(ak, sk)}"})
                    sresp.raise_for_status()
                except Exception as pe:
                    poll_fail += 1
                    print(f"[Kling] 轮询网络错误({poll_fail}/6): {pe} — 任务仍在进行，继续等", flush=True)
                    if cb: cb(f"Kling {waited}s… 网络重试{poll_fail}")
                    if poll_fail >= 6:
                        print("[Kling] 连续轮询失败过多，放弃此页", flush=True)
                        return False
                    continue
                poll_fail = 0
                jd = (sresp.json().get("data") or {})
                st = jd.get("task_status")
                if cb: cb(f"Kling {waited}s… {st}")
                if st == "succeed":
                    vids = (jd.get("task_result") or {}).get("videos") or []
                    if not vids or not vids[0].get("url"):
                        return False
                    for attempt in (1, 2):
                        try:
                            with hc.stream("GET", vids[0]["url"]) as vr:
                                vr.raise_for_status()
                                with open(out, "wb") as f:
                                    for chunk in vr.iter_bytes(65536): f.write(chunk)
                            break
                        except Exception as de:
                            print(f"[Kling] 下载失败(第{attempt}次): {de}", flush=True)
                            if attempt == 2: return False
                            time.sleep(3)
                    return out.exists() and out.stat().st_size > 1024
                if st == "failed":
                    print(f"[Kling] failed: {jd.get('task_status_msg')}", flush=True); return False
            print("[Kling] 超时", flush=True); return False
    except Exception as e:
        print(f"[Kling] {e}", flush=True); return False


def run_ppt_pipeline(job_id: str, pptx_path: Path, use_kling: bool = False, use_eleven: bool = True):
    job = JOBS[job_id]
    job.update(status="running", steps=[], progress=0)
    fid = job_id[:8]
    try:
        _log(job, "🖼 解析 PPT…", 3)
        images = pptx_to_images(pptx_path, SLIDES_DIR / fid)
        notes  = pptx_notes(pptx_path)
        texts  = pptx_slide_texts(pptx_path)
        n = len(images)
        if n == 0:
            raise RuntimeError("PPT 没有可用页面")
        def _narr(i):
            t = (notes[i] if i < len(notes) else "").strip()
            if not t:
                t = (texts[i] if i < len(texts) else "").strip()  # 备注为空→回退用页面文字
            return t
        def _sub(i):
            t = _narr(i)
            return (re.split(r"[。\n!！?？]", t)[0][:18]) if t else f"第{i+1}页"
        narr_cnt = sum(1 for i in range(n) if _narr(i))
        job["story"] = {"title": pptx_path.stem,
                        "scenes": [{"id": i+1, "subtitle_zh": _sub(i)} for i in range(n)],
                        "spiritual_application": {}}
        _log(job, f"✅ 共 {n} 页，其中 {narr_cnt} 页有旁白(备注/页面文字)", 8)
        if narr_cnt == 0:
            _log(job, "⚠️ 所有页面都没有旁白文字——成片将无配音、无字幕。请在每页 PPT 的『备注』里写旁白后重试。")

        eng_kling = bool(use_kling and kling_configured())
        if use_kling and not kling_configured():
            _log(job, "⚠️ 勾选了 Kling 但未配置 KLING_ACCESS_KEY/KLING_SECRET_KEY，回退 Ken Burns")
        _log(job, f"🎬 视频引擎：{'Kling 图生视频(付费)' if eng_kling else 'Ken Burns 镜头运动(免费)'}")
        _log(job, f"🔊 配音引擎：{'ElevenLabs(优先，失败回退)' if use_eleven else 'edge-tts(免费)'}")
        composed: list[Path] = []
        for i, img in enumerate(images):
            sid  = i + 1
            base = 8 + int(i / n * 80)
            _log(job, f"🎬 第 {sid}/{n} 页…", base)
            job["cur"] = sid
            narration = _narr(i)
            subtitle  = _sub(i) if narration else ""
            if not narration:
                _log(job, f"  ⚠️ 第 {sid} 页无备注/文字 → 静音无字幕")
            aud  = AUDIO_DIR / f"{fid}_{sid:02d}.mp3"
            vid  = CLIPS_DIR / f"{fid}_{sid:02d}.mp4"
            comp = COMP_DIR  / f"{fid}_{sid:02d}.mp4"
            if narration:
                asyncio.run(tts_to_file(narration, aud, use_elevenlabs=use_eleven))
            dur = _audio_dur(aud) if aud.exists() else 0.0
            dur = max(3.0, dur + 0.6) if dur else 4.0
            made_kling = False
            if eng_kling:
                _log(job, f"  🎞 Kling 图生视频中…(较慢)")
                if generate_kling_clip(img, narration or subtitle, dur, vid,
                                       cb=lambda m: _log(job, f"  ·{m}")):
                    made_kling = True
                    if aud.exists():
                        gap = _audio_dur(aud) - _audio_dur(vid)
                        if gap > 0.3:
                            padded = vid.with_name(vid.stem + "_pad.mp4")
                            if _pad_video(vid, gap, padded) and padded.exists():
                                vid = padded
                else:
                    _log(job, f"  ⚠️ Kling 失败，回退 Ken Burns")
            if not made_kling:
                kenburns_clip(img, dur, vid, zoom_in=(i % 2 == 0))
            if aud.exists() and aud.stat().st_size > 256:
                compose_clip(vid, aud, subtitle, comp)
                composed.append(comp if (comp.exists() and comp.stat().st_size > 1024) else vid)
            else:
                composed.append(vid)
            _log(job, f"  ✅ 第 {sid} 页完成", base + 2)

        _log(job, "🔗 FFmpeg 拼接…", 90)
        final = FILM_DIR / f"{fid}_final.mp4"
        if not concat_all(composed, final):
            raise RuntimeError("FFmpeg 拼接失败")
        mb = round(final.stat().st_size / 1024 / 1024, 1)
        _log(job, f"✅ 拼接完成 {mb} MB", 95)

        r2_url = None
        try:
            _log(job, "☁️ 上传 R2…", 97)
            r2_url = upload_r2(final)
            _log(job, f"✅ {r2_url}", 99)
        except Exception as e:
            _log(job, f"⚠️ R2 跳过: {e}")

        job.update(status="done", progress=100,
                   result={"file": final.name, "r2_url": r2_url, "mb": mb, "scenes": n})
        _log(job, "🎉 完成！")
    except Exception as e:
        import traceback; traceback.print_exc()
        job.update(status="error", error=str(e))
        _log(job, f"❌ {e}")


# ══════════════════════════════════════════════════════════════════════════════
# API 端点
# ══════════════════════════════════════════════════════════════════════════════

class StartReq(BaseModel):
    story_text:    str
    anthropic_key: str = ""
    gemini_key:    str = ""
    num_scenes:    int = 25

@router.post("/api/film/start")
def api_film_start(req: StartReq):
    ak = req.anthropic_key or os.environ.get("ANTHROPIC_API_KEY","")
    gk = req.gemini_key    or os.environ.get("GEMINI_API_KEY","")
    ck = os.environ.get("GEMINI_API_CHAT_KEY","") or gk   # 拆分镜头(chat)用独立 key，未配则回退 gk
    if not gk: raise Exception("需要 Gemini API Key")
    jid = str(uuid.uuid4())
    JOBS[jid] = {"job_id":jid,"status":"queued","progress":0,
                 "steps":[],"cur":0,"story":None,"result":None,"error":None}
    threading.Thread(target=run_pipeline,
                     args=(jid, req.story_text, ak, gk, ck, req.num_scenes),
                     daemon=True).start()
    return {"job_id": jid}


@router.post("/api/film/start-ppt")
async def api_film_start_ppt(file: UploadFile = File(...), use_kling: bool = Form(False), use_elevenlabs: bool = Form(True)):
    name = (file.filename or "").lower()
    if not name.endswith((".pptx", ".ppt")):
        raise Exception("请上传 .pptx 文件")
    jid = str(uuid.uuid4())
    pptx_path = UPLOAD_DIR / f"{jid}.pptx"
    pptx_path.write_bytes(await file.read())
    JOBS[jid] = {"job_id": jid, "status": "queued", "progress": 0,
                 "steps": [], "cur": 0, "story": None, "result": None, "error": None}
    threading.Thread(target=run_ppt_pipeline, args=(jid, pptx_path, use_kling, use_elevenlabs), daemon=True).start()
    return {"job_id": jid}

@router.get("/api/film/status/{jid}")
def api_film_status(jid: str):
    j = JOBS.get(jid)
    if not j: raise Exception("Job not found")
    return j

@router.get("/api/film/sse/{jid}")
def api_film_sse(jid: str):
    def stream() -> Generator[str, None, None]:
        seen = 0
        while True:
            j = JOBS.get(jid)
            if not j: yield 'data: {"error":"not found"}\n\n'; return
            steps = j.get("steps",[])
            payload = json.dumps({
                "status": j["status"], "progress": j["progress"],
                "cur": j.get("cur",0), "new_steps": steps[seen:],
                "result": j.get("result"), "error": j.get("error"),
                "story": j.get("story") if seen == 0 and j.get("story") else None,
            }, ensure_ascii=False)
            seen = len(steps)
            yield f"data: {payload}\n\n"
            if j["status"] in ("done","error"): return
            time.sleep(2)
    return StreamingResponse(stream(), media_type="text/event-stream",
                             headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no"})

@router.get("/api/film/download/{fname}")
def api_film_download(fname: str):
    p = FILM_DIR / fname
    if not p.exists(): raise Exception("File not found")
    return FileResponse(str(p), media_type="video/mp4", filename=fname,
                        headers={"Accept-Ranges":"bytes"})

@router.get("/film-clips/{fname}")
def api_film_clip(fname: str):
    for d in [COMP_DIR, CLIPS_DIR]:
        p = d / fname
        if p.exists():
            return FileResponse(str(p), media_type="video/mp4")
    raise Exception("Clip not found")


# ══════════════════════════════════════════════════════════════════════════════
# HTML 页面
# ══════════════════════════════════════════════════════════════════════════════

_HTML = r"""<!DOCTYPE html>
<html lang="zh">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>圣经电影工作台</title>
<style>
:root{--bg:#0a0a1a;--panel:#12122a;--card:#181830;--border:rgba(255,255,255,0.07);
  --accent:#5856d6;--gold:#ffd60a;--green:#30d158;--red:#ff453a;
  --text:rgba(255,255,255,0.92);--muted:rgba(255,255,255,0.42);}
*{box-sizing:border-box;margin:0;padding:0;}
body{background:var(--bg);color:var(--text);font-family:-apple-system,"PingFang SC",sans-serif;height:100vh;overflow:hidden;}
.layout{display:grid;grid-template-columns:400px 1fr;height:100vh;}
.left{background:var(--panel);border-right:1px solid var(--border);display:flex;flex-direction:column;overflow:hidden;}
.right{display:flex;flex-direction:column;overflow:hidden;}
/* Header */
.hdr{padding:14px 18px;background:linear-gradient(135deg,rgba(88,86,214,.25),rgba(255,214,10,.07));border-bottom:1px solid var(--border);}
.hdr h1{font-size:15px;font-weight:800;color:var(--gold);}
.hdr p{font-size:11px;color:var(--muted);margin-top:2px;}
/* Left body */
.lb{flex:1;padding:14px;display:flex;flex-direction:column;gap:9px;overflow-y:auto;}
lbl{display:block;font-size:10px;color:var(--muted);text-transform:uppercase;letter-spacing:.06em;margin-bottom:3px;}
textarea,input[type=text],input[type=password],input[type=number]{
  background:rgba(255,255,255,0.04);border:1px solid var(--border);border-radius:7px;
  color:var(--text);font-family:inherit;outline:none;transition:border .15s;width:100%;}
textarea:focus,input:focus{border-color:rgba(88,86,214,.5);}
#story{height:260px;font-size:12px;line-height:1.7;padding:11px;resize:none;}
.irow{display:flex;gap:7px;}
.irow input{padding:8px 10px;font-size:12px;}
.nrow{display:flex;align-items:center;gap:8px;font-size:12px;color:var(--muted);}
.nrow input{width:64px;padding:6px 8px;text-align:center;}
.btn{padding:10px 16px;border-radius:8px;border:none;cursor:pointer;font-size:13px;font-weight:700;transition:opacity .15s;}
.btn:hover{opacity:.86;}
.btn:disabled{opacity:.38;cursor:default;}
.btn-p{background:linear-gradient(135deg,var(--accent),#7b79f0);color:#fff;width:100%;margin-top:4px;}
.btn-dl{background:var(--green);color:#000;padding:8px 14px;font-size:12px;border-radius:7px;border:none;cursor:pointer;}
/* Right */
.rhdr{padding:12px 16px;border-bottom:1px solid var(--border);display:flex;align-items:center;gap:8px;}
.rhdr h2{font-size:13px;font-weight:700;flex:1;}
.prog-wrap{padding:12px 16px;border-bottom:1px solid var(--border);}
.prog-bar-bg{height:8px;background:rgba(255,255,255,0.08);border-radius:4px;overflow:hidden;margin:8px 0 4px;}
.prog-bar{height:100%;background:linear-gradient(90deg,var(--accent),var(--gold));border-radius:4px;transition:width .6s;}
.pct{font-size:11px;color:var(--muted);}
.grid{flex:1;overflow-y:auto;padding:12px 16px;display:grid;grid-template-columns:repeat(auto-fill,minmax(220px,1fr));gap:9px;align-content:start;}
.card{background:var(--card);border:1px solid var(--border);border-radius:9px;padding:10px;font-size:11px;transition:border .2s;}
.card.act{border-color:var(--accent);background:rgba(88,86,214,.09);}
.card.done{border-color:rgba(48,209,88,.3);}
.cnum{font-size:9px;color:var(--muted);text-transform:uppercase;letter-spacing:.06em;}
.ctit{font-weight:700;margin:2px 0;font-size:12px;}
.cst{color:var(--muted);font-size:11px;margin-top:3px;}
.cst.act{color:#5ac8fa;}.cst.ok{color:var(--green);}
video{width:100%;border-radius:6px;background:#000;margin-top:6px;max-height:120px;}
.logbox{height:140px;overflow-y:auto;padding:9px 16px;font-size:11px;line-height:1.75;color:var(--muted);border-top:1px solid var(--border);}
.logbox .lok{color:var(--green);}.logbox .lerr{color:var(--red);}
.resbar{padding:11px 16px;border-top:1px solid var(--border);background:var(--panel);display:flex;align-items:center;gap:10px;flex-wrap:wrap;}
.resbar a{color:var(--gold);text-decoration:none;font-size:12px;}
.resbar a:hover{text-decoration:underline;}
</style>
</head>
<body>
<div class="layout">
<div class="left">
  <div class="hdr"><h1>🎬 圣经电影生成工作台</h1>
    <p>PPT 插画 · Ken Burns/Kling 运动 · ElevenLabs/edge-tts 配音 · FFmpeg · R2</p></div>
  <div class="lb">
    <lbl>故事板</lbl>
    <textarea id="story" placeholder="《约瑟》(Joseph)
Style: Ancient Canaan and Imperial Egypt around 1700 BC...
Main Characters: Joseph: ...
Storyboard:
* Scene 1: Jacob presents a beautiful multicolored coat...
...
* Final scene: Joseph stands with his unified family...
属灵应用旁白：约瑟的故事告诉我们：有时候神的方法和人的方法不一样。当我们愿意顺服神时，神能成就人做不到的事情。"></textarea>
    <lbl>API Keys（空则用服务器环境变量）</lbl>
    <div class="irow"><input id="ak" type="password" placeholder="Anthropic Key (sk-ant-)"/></div>
    <div class="irow"><input id="gk" type="password" placeholder="Gemini Key (AIza...)"/></div>
    <div class="nrow"><span style="flex:1">镜头数量</span><input id="ns" type="number" value="25" min="5" max="30"></div>
    <button class="btn btn-p" id="go" onclick="start()">⚡ 开始生成完整视频</button>
    <div id="jdsp" style="font-size:10px;color:var(--muted);text-align:center;margin-top:2px"></div>
    <div style="border-top:1px solid var(--border);margin:12px 0 4px;padding-top:12px">
      <lbl>或：上传 PPT 自动生成（每页插画 + Ken Burns 镜头运动 + edge-tts 配音；旁白取自每页「备注」）</lbl>
      <input id="pptfile" type="file" accept=".pptx" style="font-size:11px;color:var(--muted);padding:6px 0"/>
      <label style="display:flex;align-items:center;gap:6px;font-size:11px;color:var(--muted);margin:4px 0 6px">
        <input type="checkbox" id="usekling"/> 用 Kling 生成真实动画（付费·较慢；留空=免费 Ken Burns）
      </label>
      <label style="display:flex;align-items:center;gap:6px;font-size:11px;color:var(--muted);margin:0 0 6px">
        <input type="checkbox" id="useeleven" checked/> 用 ElevenLabs 配音（需配置 key；失败/未配自动回退 edge-tts）
      </label>
      <button class="btn btn-p" id="goppt" onclick="startPPT()" style="background:linear-gradient(135deg,#30d158,#1f9d4d)">🖼 PPT 生成视频</button>
    </div>
  </div>
</div>

<div class="right">
  <div class="rhdr">
    <h2 id="ftitle">待生成</h2>
    <span id="badge" style="font-size:11px;color:var(--muted)"></span>
  </div>
  <div class="prog-wrap" id="pw" style="display:none">
    <div style="font-size:11px;color:var(--muted)">整体进度</div>
    <div class="prog-bar-bg"><div class="prog-bar" id="pb" style="width:0%"></div></div>
    <div class="pct" id="pct">0%</div>
  </div>
  <div class="grid" id="grid">
    <div style="grid-column:1/-1;display:flex;flex-direction:column;align-items:center;justify-content:center;height:260px;gap:12px;opacity:.3">
      <div style="font-size:44px">🎞️</div><div style="font-size:13px">粘贴故事板后点击开始</div>
    </div>
  </div>
  <div class="logbox" id="log"></div>
  <div class="resbar" id="res" style="display:none"></div>
</div>
</div>

<script>
let jid=null, es=null, scenes=[], seen=0;

function start(){
  const story=document.getElementById('story').value.trim();
  if(!story) return alert('请输入故事板');
  const ak=document.getElementById('ak').value.trim();
  const gk=document.getElementById('gk').value.trim();
  const ns=+document.getElementById('ns').value||25;
  document.getElementById('go').disabled=true;
  document.getElementById('pw').style.display='';
  document.getElementById('res').style.display='none';
  document.getElementById('log').innerHTML='';
  seen=0; scenes=[];
  document.getElementById('grid').innerHTML='<div style="color:var(--muted);padding:20px;font-size:13px">初始化…</div>';
  fetch('/api/film/start',{method:'POST',headers:{'Content-Type':'application/json'},
    body:JSON.stringify({story_text:story,anthropic_key:ak,gemini_key:gk,num_scenes:ns})})
  .then(r=>r.json()).then(d=>{
    jid=d.job_id;
    document.getElementById('jdsp').textContent='Job: '+jid.slice(0,8);
    listenSSE(jid);
  }).catch(e=>{alert('启动失败: '+e);document.getElementById('go').disabled=false;});
}

function startPPT(){
  const f=document.getElementById('pptfile').files[0];
  if(!f) return alert('请选择 .pptx 文件');
  document.getElementById('goppt').disabled=true;
  document.getElementById('go').disabled=true;
  document.getElementById('pw').style.display='';
  document.getElementById('res').style.display='none';
  document.getElementById('log').innerHTML='';
  seen=0; scenes=[];
  document.getElementById('grid').innerHTML='<div style="color:var(--muted);padding:20px;font-size:13px">上传 PPT 中…</div>';
  const fd=new FormData(); fd.append('file',f); fd.append('use_kling', document.getElementById('usekling').checked?'true':'false'); fd.append('use_elevenlabs', document.getElementById('useeleven').checked?'true':'false');
  fetch('/api/film/start-ppt',{method:'POST',body:fd})
  .then(r=>r.json()).then(d=>{
    if(d.error||!d.job_id) throw new Error(d.error||'启动失败');
    jid=d.job_id;
    document.getElementById('jdsp').textContent='Job: '+jid.slice(0,8);
    listenSSE(jid);
  }).catch(e=>{alert('启动失败: '+e);document.getElementById('goppt').disabled=false;document.getElementById('go').disabled=false;});
}

function listenSSE(id){
  if(es) es.close();
  es=new EventSource('/api/film/sse/'+id);
  es.onmessage=e=>{
    const d=JSON.parse(e.data);
    setProgress(d.progress||0);
    if(d.story && scenes.length===0) buildCards(d.story);
    if(d.new_steps) d.new_steps.forEach(addLog);
    if(d.cur>0) markScene(d.cur,'act','生成中…');
    if(d.status==='done'){onDone(d.result);es.close();}
    if(d.status==='error'){onErr(d.error);es.close();}
  };
}

function buildCards(story){
  scenes=story.scenes||[];
  document.getElementById('ftitle').textContent=story.title||'生成中…';
  const sp=story.spiritual_application||{};
  const all=[...scenes,{id:scenes.length+1,subtitle_zh:sp.title_zh||'属灵应用',_sp:true}];
  document.getElementById('grid').innerHTML=all.map(s=>`
    <div class="card" id="c${s.id}">
      <div class="cnum">${s._sp?'🙏 结尾':'Scene '+s.id}</div>
      <div class="ctit">${trunc(s.subtitle_zh||'',20)}</div>
      <div class="cst" id="cs${s.id}">待生成</div>
    </div>`).join('');
}

function markScene(id,cls,txt){
  document.querySelectorAll('.card').forEach(c=>c.classList.remove('act'));
  const c=document.getElementById('c'+id);
  const s=document.getElementById('cs'+id);
  if(c){c.classList.add(cls);c.scrollIntoView({behavior:'smooth',block:'nearest'});}
  if(s){s.className='cst '+cls;s.textContent=txt;}
}

function setProgress(p){
  document.getElementById('pb').style.width=p+'%';
  document.getElementById('pct').textContent=p+'%';
}

function addLog(msg){
  const b=document.getElementById('log');
  const p=document.createElement('p');
  p.textContent=msg;
  if(msg.includes('✅')||msg.includes('🎉')) p.className='lok';
  if(msg.includes('❌')||msg.includes('⚠️')) p.className='lerr';
  b.appendChild(p); b.scrollTop=b.scrollHeight;
}

function onDone(r){
  setProgress(100);
  document.getElementById('badge').textContent='✅ 完成';
  document.getElementById('go').disabled=false;
  document.getElementById('goppt').disabled=false;
  document.querySelectorAll('.card').forEach(c=>{c.classList.remove('act');c.classList.add('done');});
  document.querySelectorAll('[id^=cs]').forEach(e=>{e.className='cst ok';e.textContent='✅';});
  const bar=document.getElementById('res');
  bar.style.display='flex';
  let h=`<span style="color:var(--green);font-weight:700">🎉 完成 · ${r.scenes}镜头 · ${r.mb}MB</span>`;
  if(r.r2_url) h+=`<a href="${r.r2_url}" target="_blank">☁️ Cloudflare播放</a>`;
  if(r.file) h+=`<a href="/api/film/download/${r.file}" class="btn-dl" download>⬇ 下载视频</a>`;
  bar.innerHTML=h;
  addLog('🎉 全部完成！');
}
function onErr(e){
  document.getElementById('badge').textContent='❌ 错误';
  document.getElementById('go').disabled=false;
  document.getElementById('goppt').disabled=false;
  addLog('❌ '+e);
}
function trunc(s,n){return s.length>n?s.slice(0,n)+'…':s;}

window.onload=()=>{
  ['ak','gk'].forEach(id=>{
    const v=localStorage.getItem('film_'+id);
    if(v) document.getElementById(id).value=v;
    document.getElementById(id).addEventListener('change',e=>localStorage.setItem('film_'+id,e.target.value));
  });
};
</script>
</body></html>
"""

@router.get("/film-studio", response_class=HTMLResponse)
def film_studio_page():
    return _HTML
